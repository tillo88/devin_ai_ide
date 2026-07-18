"""
document_extract.py - Estrazione testo da documenti comuni allegati in chat.

Aggiunto 2026-07-09 su richiesta: "ogni tipo di documento comune leggibile" (dopo
la rimozione del vision da DEVIN — vedi local_model_launcher.py). I formati testo
puro (.txt/.md/.py/.json/.csv/.log) restano gestiti client-side in chat.html,
invariato — qui solo i formati binari che richiedono un parser vero:
PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx).

Nota onesta sui limiti:
- PDF: solo testo incorporato nel file. Un PDF scansionato (immagine pura, senza
  layer di testo) non produce nulla — servirebbe OCR, non implementato qui.
- .doc/.xls/.ppt (formati binari legacy pre-2007) NON sono supportati, solo i
  formati XML moderni (.docx/.xlsx/.pptx). python-docx/openpyxl/python-pptx non
  leggono i vecchi formati OLE.

Ogni funzione ritorna testo semplice, troncato a MAX_CHARS_PER_DOC (stessa logica
protettiva di context_engine.py per il contesto locale limitato). Nessuna eccezione
propagata all'endpoint: un errore di parsing torna come stringa leggibile nel
messaggio, cosi' la chat continua invece di rispondere con un 500.
"""

import io
from pathlib import Path

MAX_CHARS_PER_DOC = 15000

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}


def extract_text(filename: str, content: bytes) -> str:
    """Dispatcher per estensione. Ritorna testo estratto o un messaggio d'errore
    leggibile (mai un'eccezione) — pensato per finire diretto nel messaggio chat."""
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".pdf":
            text = _extract_pdf(content)
        elif ext == ".docx":
            text = _extract_docx(content)
        elif ext == ".xlsx":
            text = _extract_xlsx(content)
        elif ext == ".pptx":
            text = _extract_pptx(content)
        else:
            supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
            return (f"[Formato '{ext or 'sconosciuto'}' non supportato per l'estrazione. "
                     f"Supportati: {supported} (oltre a .txt/.md/.py/.json/.csv/.log)]")
    except Exception as e:
        return f"[Errore estrazione '{filename}': {e}]"

    text = (text or "").strip()
    if not text:
        return (f"[{filename}: nessun testo estraibile — probabile PDF scansionato "
                f"(immagine senza layer di testo, serve OCR non implementato) o file vuoto]")
    if len(text) > MAX_CHARS_PER_DOC:
        total = len(text)
        text = text[:MAX_CHARS_PER_DOC] + f"\n...[troncato, {total} caratteri totali nel documento]"
    return text


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _extract_docx(content: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"--- Foglio: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            if any(c is not None for c in row):
                parts.append(" | ".join("" if c is None else str(c) for c in row))
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs)
                if text.strip():
                    parts.append(text)
    return "\n".join(parts)
